from pptx import Presentation

def extract_text_from_pptx(file_path):
    presentation = Presentation(file_path)

    all_text = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        all_text.append(run.text)

    return all_text


file_path = "test.pptx"
texts = extract_text_from_pptx(file_path)

for text in texts:
    print(text)
